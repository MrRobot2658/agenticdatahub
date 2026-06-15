#!/usr/bin/env python3
"""把 ppt/images/slide-*.png 拼成全幅 16:9 PPT → docs/AgenticDataHub-产品介绍-v2.pptx

配合 scripts/shoot_ppt.mjs（Playwright 截图）使用，最大程度还原 ppt/index.html 的视觉效果。
用法：node scripts/shoot_ppt.mjs && .venv/bin/python scripts/build_ppt_from_shots.py
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

ROOT = Path(__file__).resolve().parent.parent
IMG_DIR = ROOT / "ppt" / "images"
OUT = ROOT / "docs" / "AgenticDataHub-产品介绍-v2.pptx"

prs = Presentation()
prs.slide_width = Emu(12192000)   # 16:9
prs.slide_height = Emu(6858000)
blank = prs.slide_layouts[6]

shots = sorted(IMG_DIR.glob("slide-*.png"))
if not shots:
    raise SystemExit("没有截图，请先运行 node scripts/shoot_ppt.mjs")

for img in shots:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(img), 0, 0, width=prs.slide_width, height=prs.slide_height)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(OUT))
print(f"已生成 {OUT}（{len(shots)} 页，全幅还原）")
